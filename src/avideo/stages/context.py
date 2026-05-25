"""ContextStage — real text extraction from .pdf / .pptx / .md documents.

Replaces ContextStub from stages/stubs.py (Phase 1) with genuine extraction.
The stage_name is kept as "context" to preserve the workdir checkpoint contract.

Supported formats:
    .pdf      -- PyMuPDF (import fitz per CLAUDE.md); raises on encrypted PDFs.
    .pptx     -- python-pptx; extracts shape text + notes (with guard).
    .md       -- plain read_text; no transformation.
    .markdown -- same as .md.

Security:
    - Suffix validated against _DISPATCH allow-list before file is opened (T-02-01).
    - Extracted content is never logged at INFO (T-02-04).
    - Token cap truncates oversized context (T-02-02, D-04).
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING

from rich.console import Console

from avideo.models.context import ContextOutput
from avideo.stages.base import CheckpointMixin

if TYPE_CHECKING:
    from avideo.models.config import RunConfig
    from avideo.utils.workdir import WorkdirManager

logger = logging.getLogger(__name__)
_console = Console(stderr=True)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

#: Maximum tokens of context text injected into the storyboard prompt (D-04).
#: ~4 chars/token (A1). Exposed as a constant so tests can derive cap_chars = N * 4.
CONTEXT_TOKEN_CAP: int = 6000


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def truncate_to_tokens(text: str, max_tokens: int = CONTEXT_TOKEN_CAP) -> str:
    """Truncate *text* to approximately *max_tokens* tokens.

    Uses the 4-chars-per-token heuristic (A1, RESEARCH).

    Args:
        text: Raw extracted text.
        max_tokens: Maximum number of tokens to keep.

    Returns:
        The original text if within budget, otherwise the first ``max_tokens * 4``
        characters.
    """
    max_chars = max_tokens * 4
    return text if len(text) <= max_chars else text[:max_chars]


def extract_pdf(path: Path) -> str:
    """Extract plain text from a PDF using PyMuPDF (fitz).

    Args:
        path: Path to the PDF file.

    Returns:
        Concatenated text from all pages.

    Raises:
        ValueError: If the PDF is password-protected.
    """
    import fitz  # noqa: PLC0415 — PyMuPDF; CLAUDE.md mandates `import fitz`

    doc = fitz.open(str(path))
    if doc.needs_pass:
        doc.close()
        raise ValueError(
            f"PDF is password-protected: {path}. "
            "Provide a decrypted copy or use a different context document."
        )
    text = "\n".join(page.get_text("text") for page in doc)
    doc.close()
    return text


def extract_pptx(path: Path) -> str:
    """Extract plain text from a .pptx file using python-pptx.

    Iterates all slide shapes for text frames, then appends notes text ONLY
    when ``slide.has_notes_slide`` is True (Pitfall 6 — reading notes_slide
    unconditionally creates an empty notes slide as a side effect).

    Args:
        path: Path to the .pptx file.

    Returns:
        Joined text of all non-empty shape texts and notes.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        # Guard: has_notes_slide BEFORE accessing notes_slide (Pitfall 6)
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(p for p in parts if p.strip())


def extract_md(path: Path) -> str:
    """Read plain text from a Markdown file.

    Args:
        path: Path to a .md or .markdown file.

    Returns:
        File contents as a UTF-8 string.
    """
    return path.read_text(encoding="utf-8")


# ---------------------------------------------------------------------------
# Dispatch map — suffix allow-list (T-02-01)
# ---------------------------------------------------------------------------

_DISPATCH: dict[str, object] = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".md": extract_md,
    ".markdown": extract_md,
}

_ALLOWED_SUFFIXES = ", ".join(sorted(_DISPATCH))


# ---------------------------------------------------------------------------
# Stage
# ---------------------------------------------------------------------------


class ContextStage(CheckpointMixin):
    """Real context ingestion stage replacing Phase-1 ContextStub.

    Extracts text from an optional .pdf / .pptx / .md context document and
    returns a ContextOutput suitable for injection into the storyboard prompt.

    stage_name is "context" — identical to ContextStub — so existing workdir
    checkpoint contracts and the orchestrator loop are unchanged.
    """

    stage_name: str = "context"

    def run(self, workdir: "WorkdirManager", config: "RunConfig") -> ContextOutput:
        """Extract context text or return ContextOutput(used=False) when absent.

        Args:
            workdir: WorkdirManager (not used for I/O in this stage; present for
                StageProtocol compliance — the orchestrator owns checkpoint writes).
            config: RunConfig; reads ``config.context`` (Optional[Path]).

        Returns:
            ContextOutput with:
            - ``used=False`` when ``config.context`` is None.
            - ``used=False`` (+ warning) when extraction yields empty/whitespace text.
            - ``used=True`` with truncated text on success.

        Raises:
            ValueError: If the suffix is not in the allow-list, or if the PDF is
                password-protected.
        """
        # CTX-02: no context document → skip cleanly
        if config.context is None:
            return ContextOutput(used=False)

        context_path: Path = config.context
        suffix = context_path.suffix.lower()

        # T-02-01: validate suffix against allow-list before opening the file
        if suffix not in _DISPATCH:
            raise ValueError(
                f"Unsupported context file type '{suffix}'. "
                f"Allowed: {_ALLOWED_SUFFIXES}"
            )

        extractor = _DISPATCH[suffix]  # type: ignore[call-arg]
        raw_text: str = extractor(context_path)  # type: ignore[operator]

        # Truncate to token cap (D-04, T-02-02)
        truncated = truncate_to_tokens(raw_text)

        # CTX-02 semantics: empty extraction → behave as if no context (Pitfall 5)
        if not truncated.strip():
            _console.print(
                f"[yellow]Warning:[/yellow] context document '{context_path.name}' "
                "yielded no extractable text — continuing without context.",
                highlight=False,
            )
            logger.debug(
                "Context extraction returned empty text for suffix=%s path=%s",
                suffix,
                context_path.name,  # never log full path or content at INFO (T-02-04)
            )
            return ContextOutput(used=False)

        logger.debug(
            "Context extracted: suffix=%s chars=%d (capped at %d tokens)",
            suffix,
            len(truncated),
            CONTEXT_TOKEN_CAP,
            # NOTE: never log the actual text content (T-02-04)
        )

        return ContextOutput(
            source_path=str(context_path),
            text=truncated,
            used=True,
        )
