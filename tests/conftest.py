"""Shared pytest fixtures for avideo test suite."""
import base64
import shutil
import subprocess
import types
import pytest
from pathlib import Path


@pytest.fixture
def tmp_workdir(tmp_path: Path) -> Path:
    """Return a workdir path inside tmp_path (not yet created)."""
    return tmp_path / "workdir"


@pytest.fixture
def minimal_bullets(tmp_path: Path) -> Path:
    """Write a minimal bullets.yaml and return its Path."""
    bullets = tmp_path / "bullets.yaml"
    bullets.write_text(
        "title: Test\nbullets:\n  - Point 1\n  - Point 2\n",
        encoding="utf-8",
    )
    return bullets


@pytest.fixture
def minimal_config(tmp_path: Path) -> Path:
    """Write a minimal config.yaml and return its Path."""
    cfg = tmp_path / "config.yaml"
    cfg.write_text(
        "voice: elevenlabs\nslides_mode: auto\nlevel: 4\nwpm: 150\n",
        encoding="utf-8",
    )
    return cfg


# ---------------------------------------------------------------------------
# Phase 2 — context ingestion fixtures (lazy imports to avoid import errors
# before 'uv add anthropic PyMuPDF python-pptx' is run in Task 1)
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_md(tmp_path: Path) -> Path:
    """Write a Markdown file with known text and return its Path."""
    md = tmp_path / "context.md"
    md.write_text(
        "# Test Context\n\nThis is the MD context text for unit tests.\n",
        encoding="utf-8",
    )
    return md


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    """Build a 1-page PDF with known text using fitz (PyMuPDF) and return its Path."""
    import fitz  # noqa: PLC0415 — lazy import; PyMuPDF installed in Task 1

    pdf_path = tmp_path / "context.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "PDF context text for unit tests.")
    doc.save(str(pdf_path))
    doc.close()
    return pdf_path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Build a 1-slide .pptx with known text using python-pptx and return its Path."""
    from pptx import Presentation  # noqa: PLC0415 — lazy import; python-pptx installed in Task 1
    from pptx.util import Inches, Pt  # noqa: PLC0415

    pptx_path = tmp_path / "context.pptx"
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.text = "PPTX context text for unit tests."
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def encrypted_pdf(tmp_path: Path) -> Path:
    """Build a password-protected PDF using fitz and return its Path."""
    import fitz  # noqa: PLC0415 — lazy import; PyMuPDF installed in Task 1

    pdf_path = tmp_path / "encrypted.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Secret content.")
    doc.save(
        str(pdf_path),
        encryption=fitz.PDF_ENCRYPT_AES_256,
        owner_pw="owner",
        user_pw="user",
    )
    doc.close()
    return pdf_path


# ---------------------------------------------------------------------------
# Phase 4 — voice / alignment fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def fake_elevenlabs_response():
    """Return a minimal ElevenLabs convert_with_timestamps response object.

    Attributes mirror CharacterAlignmentResponseModel in elevenlabs SDK 2.x
    (character_start_times_seconds / character_end_times_seconds in SECONDS —
    NOT the obsolete _ms fields of SDK 1.x).

    The timestamps are strictly increasing so the default behaviour is success.
    Characters cover "hola mundo" (10 chars incl. space).
    """
    alignment = types.SimpleNamespace(
        characters=["h", "o", "l", "a", " ", "m", "u", "n", "d", "o"],
        character_start_times_seconds=[0.00, 0.05, 0.10, 0.15, 0.20, 0.25, 0.30, 0.35, 0.40, 0.45],
        character_end_times_seconds=[0.05, 0.10, 0.15, 0.20, 0.25, 0.30, 0.35, 0.40, 0.45, 0.50],
    )
    # A few bytes of valid base64-encoded content (not a real mp3, but decodeable)
    audio_bytes = b"\xff\xe3\x10\x00"  # minimal mp3-like header bytes
    return types.SimpleNamespace(
        audio_base64=base64.b64encode(audio_bytes).decode("utf-8"),
        alignment=alignment,
    )


@pytest.fixture
def fake_word_segments():
    """Return minimal whisperx-style word_segments list for 'hola mundo'.

    Used by plan 04-02 (align stage) tests.  start/end are SECONDS relative to
    the beginning of the slide clip.
    """
    return [
        {"word": "hola", "start": 0.0, "end": 0.4},
        {"word": "mundo", "start": 0.5, "end": 0.9},
    ]


@pytest.fixture
def voice_config(tmp_path: Path):
    """Return a minimal RunConfig for voice tests using tmp_path as workdir.

    Mirrors the _build_config pattern used in test_storyboard.py.
    Sets voice=elevenlabs, voice_id to a test placeholder, and points the
    workdir at a tmp_path subdirectory.
    """
    from avideo.models.config import RunConfig, VoiceMode  # noqa: PLC0415 — lazy

    bullets = tmp_path / "bullets.yaml"
    bullets.write_text("title: Test\nbullets:\n  - Point 1\n", encoding="utf-8")
    return RunConfig(
        bullets=bullets,
        duration=60,
        voice=VoiceMode.elevenlabs,
        voice_id="test-voice-id",
        workdir=tmp_path / "workdir",
    )


# ---------------------------------------------------------------------------
# Phase 3 — slides rendering fixtures (lazy imports; deps added in 03-01 Task 1)
# ---------------------------------------------------------------------------


@pytest.fixture
def fake_storyboard():
    """Return a StoryboardOutput with one slide for EACH of the 7 VisualType values.

    Covers every macro-dispatch path so plan 03-02 tests can assert that all
    visual_type values render without KeyError.  Chart/comparison slides include
    bullets with numbers/percentages to exercise the numeric-parsing path.
    """
    from avideo.models.storyboard import (  # noqa: PLC0415 — lazy
        SlideSpec,
        StoryboardOutput,
        VisualType,
    )

    return StoryboardOutput(
        language="es",
        slides=[
            SlideSpec(
                title="Bienvenidos al pipeline de vídeo",
                bullets=["Automatización completa", "Sin edición manual"],
                visual_type=VisualType.title,
            ),
            SlideSpec(
                title="Características principales",
                bullets=[
                    "Slides generadas automáticamente",
                    "Voz en off sincronizada",
                    "Subtítulos precisos",
                ],
                visual_type=VisualType.bullets,
            ),
            SlideSpec(
                title="Crecimiento de adopción",
                bullets=[
                    "Ventas Q1: 40%",
                    "Coste reducción: 25%",
                    "Productividad +35%",
                ],
                visual_type=VisualType.chart,
            ),
            SlideSpec(
                title="Flujo del pipeline",
                bullets=[
                    "Bullets → Storyboard",
                    "Storyboard → Slides",
                    "Slides + Voz → Vídeo",
                ],
                visual_type=VisualType.diagram,
            ),
            SlideSpec(
                title="La automatización no reemplaza la creatividad",
                bullets=[
                    "Libera tiempo para pensar",
                    "— Equipo de producto",
                ],
                visual_type=VisualType.quote,
            ),
            SlideSpec(
                title="Antes vs. Después",
                bullets=[
                    "Manual: 4 horas por vídeo",
                    "Automático: 10 minutos por vídeo",
                    "Ahorro: 95%",
                ],
                visual_type=VisualType.comparison,
            ),
            SlideSpec(
                title="Tecnologías utilizadas",
                bullets=[
                    "Playwright para render",
                    "ElevenLabs para voz",
                    "FFmpeg para montaje",
                ],
                visual_type=VisualType.image_icon,
            ),
        ],
    )


# ---------------------------------------------------------------------------
# Phase 5 — assembly / FFmpeg fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def loudnorm_pass1_stderr():
    """Return a canned ffmpeg stderr string containing the loudnorm pass-1 JSON block.

    Field names are VERIFIED from ffmpeg 8.0.1 loudnorm output (05-RESEARCH Pattern 5).
    Used by plan 05-02 parse_loudnorm_json tests.
    """
    # Simulate realistic loudnorm stderr with log lines surrounding the JSON block
    return (
        "ffmpeg version 8.0.1\n"
        "[Parsed_loudnorm_0 @ 0x...] Input Integrated:    -22.01 LUFS\n"
        "[Parsed_loudnorm_0 @ 0x...] Input True Peak:    -20.91 dBTP\n"
        "[Parsed_loudnorm_0 @ 0x...] Input LRA:           0.70 LU\n"
        "[Parsed_loudnorm_0 @ 0x...] Input Threshold:   -32.01 LUFS\n"
        "[Parsed_loudnorm_0 @ 0x...]\n"
        "{\n"
        '    "input_i" : "-22.01",\n'
        '    "input_tp" : "-20.91",\n'
        '    "input_lra" : "0.70",\n'
        '    "input_thresh" : "-32.01",\n'
        '    "output_i" : "-15.74",\n'
        '    "output_tp" : "-14.60",\n'
        '    "output_lra" : "0.50",\n'
        '    "output_thresh" : "-25.74",\n'
        '    "normalization_type" : "dynamic",\n'
        '    "target_offset" : "-0.26"\n'
        "}\n"
    )


@pytest.fixture
def tiny_av_assets(tmp_path: Path):
    """Generate N tiny PNGs and N short sine-wave audio files for smoke tests.

    Uses Pillow for PNGs and ffmpeg lavfi for audio (sine=frequency=440:duration=1).
    Skips silently (returns empty lists) when ffmpeg is absent — the smoke test
    that uses this fixture is itself guarded by a skipif.

    Returns:
        tuple[list[str], list[str]]: (png_paths, audio_paths)
            Both lists have the same length (N=3).
    """
    from PIL import Image  # noqa: PLC0415 — Pillow is a dev dep

    n = 3
    png_paths = []
    audio_paths = []

    # Create tiny PNG files (red/green/blue 32x32 images)
    colors = [(220, 50, 50), (50, 200, 50), (50, 50, 220)]
    for i, color in enumerate(colors[:n]):
        img = Image.new("RGB", (32, 32), color=color)
        png_path = tmp_path / f"slide_{i:02d}.png"
        img.save(str(png_path))
        png_paths.append(str(png_path))

    # Create short audio files via ffmpeg lavfi (sine wave, ~1s each)
    if shutil.which("ffmpeg") is None:
        # Return dummy paths — the guarded smoke test will skip before using them
        return (png_paths, [str(tmp_path / f"audio_{i:02d}.mp3") for i in range(n)])

    for i in range(n):
        audio_path = tmp_path / f"audio_{i:02d}.mp3"
        # Duration ~1.0s per audio; slightly varied to test non-uniform durations
        dur = 1.0 + i * 0.1  # 1.0, 1.1, 1.2 seconds
        cmd = [
            "ffmpeg", "-hide_banner", "-loglevel", "error",
            "-f", "lavfi",
            "-i", f"sine=frequency=440:duration={dur}",
            "-y",
            str(audio_path),
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        audio_paths.append(str(audio_path))

    return (png_paths, audio_paths)
